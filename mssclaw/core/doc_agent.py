# -*- coding: utf-8 -*-
"""
S-008: Doc-Agent — Document pipeline for MSSclaw

Capabilities:
  doc_export   — structured data → .docx (tables, formatted text)
  doc_import   — .docx → structured data (text, tables, metadata)
  xlsx_read    — .xlsx → 2D arrays + schema inference
  xlsx_write   — structured data → .xlsx with formatting
  pdf_generate — markdown/text → PDF (via fpdf2, proven CJK)
  ppt_generate — outline → .pptx (title + bullet slides)

Design:
  - Same pattern as other S-007 agents
  - Registers with AGENT_REGISTRY in specialized_agents.py
  - Uses MeetingRoom for document state tracking
  - Validated by NormativeField before file I/O
"""
import json
import os
import uuid
import re
from typing import Dict, List, Optional, Any
from dataclasses import dataclass, field
from datetime import datetime

from docx import Document
from docx.shared import Pt, Inches, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt

# Swarm imports (fallback stubs for standalone testing)
try:
    from meeting_room import MeetingRoom
except ImportError:
    MeetingRoom = object


# ═══════════════════════════════════════════════════════
# Data classes
# ═══════════════════════════════════════════════════════

@dataclass
class DocExportResult:
    path: str
    size_kb: float
    sections: int
    tables: int
    paragraphs: int

@dataclass
class DocImportResult:
    path: str
    text: str
    tables: List[List[List[str]]]
    metadata: Dict[str, str]
    paragraph_count: int
    table_count: int

@dataclass
class XlsxReadResult:
    path: str
    sheets: Dict[str, List[List[Any]]]
    schema: Dict[str, Dict[str, str]]  # col → type
    row_counts: Dict[str, int]

@dataclass
class XlsxWriteResult:
    path: str
    sheets: int
    total_rows: int
    size_kb: float

@dataclass
class PdfResult:
    path: str
    pages: int
    size_kb: float

@dataclass
class PptResult:
    path: str
    slides: int
    size_kb: float


# ═══════════════════════════════════════════════════════
# Doc-Agent
# ═══════════════════════════════════════════════════════

class DocAgent:
    """Document pipeline agent: import/export across .docx, .xlsx, .pdf, .pptx."""
    
    CAPABILITIES = ["doc_export", "doc_import", "xlsx_read", "xlsx_write",
                    "pdf_generate", "ppt_generate"]
    
    def __init__(self, agent_id: str, room=None, swarm=None):
        self.agent_id = agent_id
        self.room = room
        self.swarm = swarm
        self._documents: Dict[str, str] = {}  # doc_id → path
        self._export_history: List[Dict] = []
    
    # ── .docx Export ──────────────────────────────────
    
    def doc_export(self, data: Dict, output_path: str) -> DocExportResult:
        """
        Export structured data to .docx.
        
        data = {
            "title": "Experiment Report",
            "author": "MSS-AI",
            "sections": [
                {"heading": "Abstract", "body": "..."},
                {"heading": "Results", "body": "...", "table": {
                    "headers": ["Model", "Eta", "Breach"],
                    "rows": [["qwen7b", "0.74", "No"], ["mss-ai", "0.77", "No"]]
                }},
            ]
        }
        """
        doc = Document()
        
        # Metadata
        core = doc.core_properties
        core.title = data.get("title", "Untitled")
        core.author = data.get("author", "MSS-AI")
        core.created = datetime.now()
        
        # Title page
        title = doc.add_heading(data.get("title", "Untitled"), level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        if data.get("subtitle"):
            sub = doc.add_paragraph(data["subtitle"])
            sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
            sub.runs[0].font.size = Pt(14)
            sub.runs[0].font.color.rgb = RGBColor(100, 100, 100)
        
        if data.get("date"):
            d = doc.add_paragraph(data["date"])
            d.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        doc.add_page_break()
        
        table_count = 0
        
        for section in data.get("sections", []):
            if section.get("heading"):
                doc.add_heading(section["heading"], level=1)
            
            if section.get("body"):
                para = doc.add_paragraph(section["body"])
                para.style.font.size = Pt(11)
            
            if section.get("table"):
                headers = section["table"]["headers"]
                rows = section["table"]["rows"]
                t = doc.add_table(rows=1 + len(rows), cols=len(headers))
                t.style = 'Light Grid Accent 1'
                t.alignment = WD_TABLE_ALIGNMENT.CENTER
                
                # Header row
                for j, h in enumerate(headers):
                    cell = t.rows[0].cells[j]
                    cell.text = str(h)
                    for run in cell.paragraphs[0].runs:
                        run.bold = True
                
                # Data rows
                for i, row in enumerate(rows):
                    for j, val in enumerate(row):
                        t.rows[i + 1].cells[j].text = str(val)
                
                doc.add_paragraph()  # spacer
                table_count += 1
        
        os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
        doc.save(output_path)
        
        size_kb = os.path.getsize(output_path) / 1024
        result = DocExportResult(
            path=output_path, size_kb=size_kb,
            sections=len(data.get("sections", [])),
            tables=table_count,
            paragraphs=sum(1 for s in data.get("sections", []) if s.get("body"))
        )
        self._export_history.append({"type": "docx", "result": result.__dict__, "ts": datetime.now().isoformat()})
        return result
    
    # ── .docx Import ──────────────────────────────────
    
    def doc_import(self, path: str) -> DocImportResult:
        """Read .docx and extract text + tables."""
        doc = Document(path)
        
        # Extract all text
        all_text = []
        for para in doc.paragraphs:
            if para.text.strip():
                all_text.append(para.text)
        
        # Extract tables
        tables = []
        for t in doc.tables:
            table_data = []
            for row in t.rows:
                table_data.append([cell.text for cell in row.cells])
            tables.append(table_data)
        
        # Metadata
        core = doc.core_properties
        metadata = {
            "title": core.title or "",
            "author": core.author or "",
            "created": str(core.created) if core.created else "",
            "modified": str(core.modified) if core.modified else "",
            "paragraphs": str(len(doc.paragraphs)),
            "tables": str(len(doc.tables)),
        }
        
        return DocImportResult(
            path=path,
            text="\n\n".join(all_text),
            tables=tables,
            metadata=metadata,
            paragraph_count=len(doc.paragraphs),
            table_count=len(doc.tables)
        )
    
    # ── .xlsx Read ────────────────────────────────────
    
    def xlsx_read(self, path: str, sheet_names: Optional[List[str]] = None) -> XlsxReadResult:
        """Read .xlsx into structured data with type inference."""
        wb = openpyxl.load_workbook(path, data_only=True)
        
        sheets_to_read = sheet_names or wb.sheetnames
        sheets = {}
        schema = {}
        row_counts = {}
        
        for name in sheets_to_read:
            if name not in wb.sheetnames:
                continue
            
            ws = wb[name]
            raw = []
            for row in ws.iter_rows(values_only=True):
                raw.append(list(row))
            
            # Remove fully empty trailing rows
            while raw and all(c is None for c in raw[-1]):
                raw.pop()
            
            sheets[name] = raw
            row_counts[name] = len(raw)
            
            # Type inference from first 5 data rows
            if len(raw) > 1:
                col_types = {}
                for j in range(len(raw[0])):
                    vals = [raw[i][j] for i in range(1, min(len(raw), 6)) if raw[i][j] is not None]
                    col_types[str(j)] = self._infer_type(vals)
                schema[name] = col_types
        
        wb.close()
        return XlsxReadResult(path=path, sheets=sheets, schema=schema, row_counts=row_counts)
    
    def _infer_type(self, values: List) -> str:
        if not values:
            return "empty"
        types = set()
        for v in values:
            if isinstance(v, (int, float)):
                types.add("number")
            elif isinstance(v, bool):
                types.add("boolean")
            elif isinstance(v, str):
                if re.match(r'^\d{4}-\d{2}-\d{2}', v):
                    types.add("date")
                elif re.match(r'^[\d,.]+$', v):
                    types.add("number")
                else:
                    types.add("string")
        return "/".join(sorted(types)) if types else "empty"
    
    # ── .xlsx Write ───────────────────────────────────
    
    def xlsx_write(self, data: Dict, output_path: str) -> XlsxWriteResult:
        """
        Write structured data to .xlsx.
        
        data = {
            "sheets": {
                "E4 Results": {
                    "headers": ["Model", "Baseline", "With Guard"],
                    "rows": [["qwen7b", 0.687, 0.952], ["mss-ai", 0.771, 0.859]],
                    "formats": {"header_bg": "4472C4", "header_font": "FFFFFF"}
                }
            }
        }
        """
        wb = openpyxl.Workbook()
        # Remove default sheet
        wb.remove(wb.active)
        
        header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
        header_font = Font(color="FFFFFF", bold=True, size=11)
        data_font = Font(size=10)
        thin_border = Border(
            left=Side(style='thin', color='D0D0D0'),
            right=Side(style='thin', color='D0D0D0'),
            top=Side(style='thin', color='D0D0D0'),
            bottom=Side(style='thin', color='D0D0D0'),
        )
        
        total_rows = 0
        
        for sheet_name, sheet_data in data.get("sheets", {}).items():
            ws = wb.create_sheet(title=sheet_name[:31])  # Excel 31-char limit
            
            headers = sheet_data.get("headers", [])
            rows = sheet_data.get("rows", [])
            
            # Write headers
            for j, h in enumerate(headers):
                cell = ws.cell(row=1, column=j + 1, value=h)
                cell.fill = header_fill
                cell.font = header_font
                cell.alignment = Alignment(horizontal='center')
                cell.border = thin_border
            
            # Write data
            for i, row in enumerate(rows):
                for j, val in enumerate(row):
                    cell = ws.cell(row=i + 2, column=j + 1, value=val)
                    cell.font = data_font
                    cell.border = thin_border
                    if isinstance(val, float):
                        cell.number_format = '0.000'
            
            # Auto-width
            for j in range(len(headers)):
                max_width = len(str(headers[j])) if j < len(headers) else 10
                for i in range(len(rows)):
                    if j < len(rows[i]):
                        max_width = max(max_width, len(str(rows[i][j])))
                ws.column_dimensions[get_column_letter(j + 1)].width = min(max_width + 4, 40)
            
            total_rows += len(rows)
        
        os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
        wb.save(output_path)
        
        size_kb = os.path.getsize(output_path) / 1024
        return XlsxWriteResult(
            path=output_path,
            sheets=len(data.get("sheets", {})),
            total_rows=total_rows,
            size_kb=size_kb
        )
    
    # ── PDF Generate ──────────────────────────────────
    
    def pdf_generate(self, content: str, output_path: str, title: str = "") -> PdfResult:
        """Render markdown/text to PDF with CJK support."""
        
        class _Pdf(FPDF):
            def __init__(self):
                super().__init__('P', 'mm', 'A4')
                self.cjk_font = None
                self.bold_font = None
                for fp, name in [
                    (r'C:\Windows\Fonts\simsun.ttc', 'CJK'),
                    (r'C:\Windows\Fonts\msyh.ttc', 'CJK'),
                    (r'C:\Windows\Fonts\msyhbd.ttc', 'CJKb'),
                ]:
                    try:
                        if 'bd' in fp.lower():
                            self.add_font('CJKb', '', fp)
                            self.bold_font = 'CJKb'
                        else:
                            self.add_font('CJK', '', fp)
                            self.cjk_font = 'CJK'
                    except:
                        pass
                self.cjk_font = self.cjk_font or 'Helvetica'
                self.bold_font = self.bold_font or self.cjk_font
                self.set_auto_page_break(True, 20)
                self.set_left_margin(20)
                self.set_right_margin(18)
                self.width = self.w - self.l_margin - self.r_margin
            
            def add_title(self, text):
                self.ln(6)
                self.set_font(self.bold_font, '', 16)
                self.multi_cell(self.width, 8, text, align='C')
                self.ln(4)
            
            def add_body(self, text):
                self.set_font(self.cjk_font, '', 10)
                paragraphs = text.split('\n\n')
                for p in paragraphs:
                    p = p.strip()
                    if not p:
                        continue
                    # Handle headers
                    if p.startswith('# ') and not p.startswith('## '):
                        self.set_font(self.bold_font, '', 14)
                        self.multi_cell(self.width, 7, p[2:])
                        self.ln(2)
                    elif p.startswith('## '):
                        self.set_font(self.bold_font, '', 12)
                        self.multi_cell(self.width, 6, p[3:])
                        self.ln(2)
                    elif p.startswith('### '):
                        self.set_font(self.bold_font, '', 11)
                        self.multi_cell(self.width, 5.5, p[4:])
                        self.ln(1)
                    elif p.startswith('```'):
                        self._code_block(p)
                    elif p.startswith('|'):
                        self._simple_table(p)
                    else:
                        self.set_font(self.cjk_font, '', 10)
                        self.multi_cell(self.width, 5.5, p)
                        self.ln(1)
            
            def _code_block(self, text):
                lines = text.replace('```', '').strip().split('\n')
                self.set_font('Courier', '', 7.5)
                self.set_fill_color(248, 248, 248)
                for l in lines:
                    self.set_x(self.l_margin + 4)
                    self.cell(self.width - 8, 4.5, l[:120], fill=True)
                    self.ln()
                self.ln(3)
            
            def _simple_table(self, text):
                rows = [[c.strip() for c in r.split('|') if c.strip()] for r in text.strip().split('\n')]
                rows = [r for r in rows if r and not all(re.match(r'^[\s\-:]+$', c) for c in r)]
                if not rows:
                    return
                n = len(rows[0])
                cw = self.width / n
                self.set_font(self.cjk_font, '', 8)
                for i, row in enumerate(rows):
                    for j, cell in enumerate(row[:n]):
                        self.cell(cw, 5.5, str(cell)[:60], border=1,
                                  fill=(i == 0), link='')
                    self.ln()
                self.ln(2)
        
        pdf = _Pdf()
        pdf.add_page()
        if title:
            pdf.add_title(title)
        pdf.add_body(content)
        
        os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
        pdf.output(output_path)
        
        return PdfResult(path=output_path, pages=pdf.page_no(),
                         size_kb=os.path.getsize(output_path) / 1024)
    
    # ── .pptx Generate ────────────────────────────────
    
    def ppt_generate(self, outline: Dict, output_path: str) -> PptResult:
        """
        Generate .pptx from outline.
        
        outline = {
            "title": "MSS-AI Quarterly Review",
            "subtitle": "Q2 2026",
            "slides": [
                {"title": "Overview", "bullets": ["Point 1", "Point 2"],
                 "note": "Speaker note text"},
                {"title": "Results", "bullets": ["Eta: 0.77", "Breach: 0%"],
                 "layout": "two_column", "col2": ["Detail A", "Detail B"]},
            ]
        }
        """
        prs = Presentation()
        prs.slide_width = PptInches(13.333)  # 16:9
        prs.slide_height = PptInches(7.5)
        
        # Title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = outline.get("title", "Untitled")
        if outline.get("subtitle"):
            title_slide.placeholders[1].text = outline["subtitle"]
        
        # Content slides
        for slide_data in outline.get("slides", []):
            slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
            
            # Title
            slide.shapes.title.text = slide_data.get("title", "")
            
            # Bullets
            if slide_data.get("bullets"):
                body = slide.placeholders[1]
                tf = body.text_frame
                tf.clear()
                for i, bullet in enumerate(slide_data["bullets"]):
                    if i == 0:
                        tf.text = bullet
                    else:
                        p = tf.add_paragraph()
                        p.text = bullet
                        p.level = 0
            
            # Speaker notes
            if slide_data.get("note"):
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = slide_data["note"]
        
        os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
        prs.save(output_path)
        
        return PptResult(path=output_path, slides=len(outline.get("slides", [])) + 1,
                         size_kb=os.path.getsize(output_path) / 1024)
    
    # ── Cloud Backup ───────────────────────────────────
    
    def cloud_backup(self, path: str) -> Dict[str, Any]:
        """
        Upload file to cloud storage (Tencent SMH).
        Calls cloud-upload-backup skill CLI.
        Returns {"ok": bool, "message": str, "url": str}.
        """
        import subprocess
        backup_cmd = os.path.expandvars(
            r"%USERPROFILE%\.qclaw\skills\cloud-upload-backup\scripts\windows\cloud_backup.cmd"
        )
        if not os.path.exists(backup_cmd):
            return {"ok": False, "message": "cloud_backup.cmd not found", "url": ""}
        try:
            result = subprocess.run(
                [backup_cmd, "upload", path],
                capture_output=True, text=True, timeout=60,
                cwd=os.path.dirname(backup_cmd)
            )
            output = result.stdout + result.stderr
            # Try to extract JSON from output
            for line in output.split('\n'):
                line = line.strip()
                if line.startswith('{'):
                    try:
                        data = json.loads(line)
                        return {"ok": True, "message": data.get("message", output[:200]),
                                "url": data.get("fileUrl", "")}
                    except:
                        pass
            return {"ok": result.returncode == 0, "message": output[:500], "url": ""}
        except Exception as e:
            return {"ok": False, "message": str(e), "url": ""}
    
    def backup_batch(self, paths: List[str]) -> List[Dict]:
        """Backup multiple files. Returns list of results."""
        return [self.cloud_backup(p) for p in paths]
    
    # ── Notification ───────────────────────────────────
    
    def notify_complete(self, title: str, details: Dict, output_dir: str = "") -> Dict:
        """
        Generate notification manifest when experiment/training completes.
        Writes notify.json for downstream delivery (email/IM/webhook).
        """
        manifest = {
            "timestamp": datetime.now().isoformat(),
            "title": title,
            "summary": details.get("summary", ""),
            "results": details.get("results", {}),
            "files": details.get("files", []),
            "agent_id": self.agent_id,
        }
        if output_dir:
            os.makedirs(output_dir, exist_ok=True)
            notify_path = os.path.join(output_dir, "notify.json")
            with open(notify_path, 'w', encoding='utf-8') as f:
                json.dump(manifest, f, ensure_ascii=False, indent=2)
            manifest["notify_path"] = notify_path
        return manifest
    
    # ── Enhanced PPT ───────────────────────────────────
    
    def ppt_from_report(self, report_data: Dict, output_path: str) -> PptResult:
        """
        Generate rich PPT from structured report data.
        Supports: data tables as charts, comparison slides, conclusion slide.
        
        report_data = {
            "title": "Q2 Review", "subtitle": "2026",
            "sections": [
                {"title": "Results", "type": "table",
                 "headers": ["Model", "Score"], "rows": [["A", 0.9]]},
                {"title": "Comparison", "type": "comparison",
                 "left": {"label": "Before", "items": ["0.6", "slow"]},
                 "right": {"label": "After", "items": ["0.95", "fast"]}},
                {"title": "Conclusion", "type": "bullets",
                 "bullets": ["Finding 1", "Finding 2"]},
            ]
        }
        """
        prs = Presentation()
        prs.slide_width = PptInches(13.333)
        prs.slide_height = PptInches(7.5)
        
        # Colors
        BLUE = (68, 114, 196)
        DARK = (50, 50, 50)
        LIGHT = (240, 240, 240)
        
        # Title slide
        ts = prs.slides.add_slide(prs.slide_layouts[0])
        ts.shapes.title.text = report_data.get("title", "Untitled")
        if report_data.get("subtitle"):
            ts.placeholders[1].text = report_data["subtitle"]
        
        for sec in report_data.get("sections", []):
            sec_type = sec.get("type", "bullets")
            
            if sec_type == "table":
                slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
                # Title
                txBox = slide.shapes.add_textbox(PptInches(0.5), PptInches(0.3),
                                                  PptInches(12), PptInches(0.8))
                tf = txBox.text_frame
                tf.text = sec.get("title", "")
                tf.paragraphs[0].font.size = PptPt(28)
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.color.rgb = RGBColor(*BLUE)
                
                # Table
                headers = sec.get("headers", [])
                rows = sec.get("rows", [])
                if headers and rows:
                    n_cols = len(headers)
                    n_rows = len(rows) + 1
                    tbl = slide.shapes.add_table(n_rows, n_cols,
                        PptInches(1), PptInches(1.5), PptInches(11), PptInches(0.5 * n_rows))
                    table = tbl.table
                    for j, h in enumerate(headers):
                        cell = table.cell(0, j)
                        cell.text = str(h)
                        for p in cell.text_frame.paragraphs:
                            p.font.bold = True
                            p.font.size = PptPt(12)
                    for i, row in enumerate(rows):
                        for j, val in enumerate(row[:n_cols]):
                            table.cell(i + 1, j).text = str(val)
            
            elif sec_type == "comparison":
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                txBox = slide.shapes.add_textbox(PptInches(0.5), PptInches(0.3),
                                                  PptInches(12), PptInches(0.8))
                tf = txBox.text_frame
                tf.text = sec.get("title", "")
                tf.paragraphs[0].font.size = PptPt(28)
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.color.rgb = RGBColor(*BLUE)
                
                left_data = sec.get("left", {})
                right_data = sec.get("right", {})
                
                # Left column
                lb = slide.shapes.add_textbox(PptInches(1), PptInches(1.5),
                                               PptInches(5), PptInches(5))
                ltf = lb.text_frame
                ltf.text = left_data.get("label", "Before")
                ltf.paragraphs[0].font.size = PptPt(20)
                ltf.paragraphs[0].font.bold = True
                for item in left_data.get("items", []):
                    p = ltf.add_paragraph()
                    p.text = f"• {item}"
                    p.font.size = PptPt(16)
                
                # Right column
                rb = slide.shapes.add_textbox(PptInches(7), PptInches(1.5),
                                               PptInches(5), PptInches(5))
                rtf = rb.text_frame
                rtf.text = right_data.get("label", "After")
                rtf.paragraphs[0].font.size = PptPt(20)
                rtf.paragraphs[0].font.bold = True
                for item in right_data.get("items", []):
                    p = rtf.add_paragraph()
                    p.text = f"• {item}"
                    p.font.size = PptPt(16)
            
            else:  # bullets
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = sec.get("title", "")
                slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(*BLUE)
                body = slide.placeholders[1]
                tf = body.text_frame
                tf.clear()
                bullets = sec.get("bullets", [])
                for i, b in enumerate(bullets):
                    if i == 0:
                        tf.text = b
                    else:
                        p = tf.add_paragraph()
                        p.text = b
                        p.font.size = PptPt(18)
        
        os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
        prs.save(output_path)
        total_slides = 1 + len(report_data.get("sections", []))
        return PptResult(path=output_path, slides=total_slides,
                         size_kb=os.path.getsize(output_path) / 1024)
    
    # ── Pipeline: experiment → report ─────────────────
    
    def experiment_to_report(self, experiment_data: Dict, output_dir: str,
                            cloud: bool = False, notify: bool = False) -> Dict[str, Any]:
        """
        One-shot: experiment JSON → .docx + .xlsx + .pdf (+ cloud + notify).
        Returns {"docx": path, "xlsx": path, "pdf": path, "cloud": [], "notify": {}}.
        """
        os.makedirs(output_dir, exist_ok=True)
        base = experiment_data.get("name", "report")
        paths = {}
        
        # .xlsx — raw data tables
        if experiment_data.get("tables"):
            xlsx_path = os.path.join(output_dir, f"{base}_data.xlsx")
            self.xlsx_write({"sheets": experiment_data["tables"]}, xlsx_path)
            paths["xlsx"] = xlsx_path
        
        # .docx — formatted report
        if experiment_data.get("report"):
            docx_path = os.path.join(output_dir, f"{base}_report.docx")
            self.doc_export(experiment_data["report"], docx_path)
            paths["docx"] = docx_path
        
        # .pdf — quick summary
        if experiment_data.get("summary_md"):
            pdf_path = os.path.join(output_dir, f"{base}_summary.pdf")
            self.pdf_generate(experiment_data["summary_md"], pdf_path,
                            title=experiment_data.get("name", ""))
            paths["pdf"] = pdf_path
        
        # Cloud backup
        cloud_results = []
        if cloud:
            for key, p in paths.items():
                cloud_results.append({"file": key, "result": self.cloud_backup(p)})
            paths["cloud"] = cloud_results
        
        # Notification
        notify_result = {}
        if notify:
            file_list = [p for p in paths.values() if isinstance(p, str)]
            notify_result = self.notify_complete(
                f"Experiment: {base}",
                {"summary": experiment_data.get("summary_md", "")[:500],
                 "results": {k: v for k, v in experiment_data.items() if k != "summary_md"},
                 "files": file_list},
                output_dir
            )
            paths["notify"] = notify_result
        
        return paths


# ═══════════════════════════════════════════════════════
# Self-test
# ═══════════════════════════════════════════════════════

def _make_test_dir():
    d = os.path.join(os.path.dirname(__file__), "_test_output")
    os.makedirs(d, exist_ok=True)
    return d


def _test():
    agent = DocAgent("doc_001")
    td = _make_test_dir()
    
    # ── T1: .docx export + import round-trip ─────
    data = {
        "title": "E4 Cross-Model Guarded Comparison",
        "subtitle": "Baseline vs GuardProxy",
        "date": "2026-06-14",
        "author": "MSS-AI",
        "sections": [
            {"heading": "Abstract", "body": "We benchmark three models under GuardProxy."},
            {"heading": "Results", "body": "qwen7b gains +0.349 under guard; mss-ai loses -0.012.",
             "table": {
                 "headers": ["Model", "Baseline η", "Guarded η", "Δ"],
                 "rows": [["qwen2.5:7b", "0.687", "0.952", "+0.349"],
                          ["mss-ai", "0.771", "0.859", "-0.012"],
                          ["qwen2.5:0.5b", "0.650", "0.650", "0.000"]]
             }},
        ]
    }
    docx_path = os.path.join(td, "e4_test.docx")
    r = agent.doc_export(data, docx_path)
    assert os.path.exists(docx_path), "docx not created"
    assert r.tables == 1
    print(f"T1 PASS: .docx export ({r.size_kb:.1f} KB, {r.tables} table)")
    
    # Round-trip
    imp = agent.doc_import(docx_path)
    assert "qwen7b" in imp.text or any("qwen7b" in str(row) for t in imp.tables for row in t)
    print(f"T2 PASS: .docx import ({imp.paragraph_count} paras, {imp.table_count} tables)")
    
    # ── T3: .xlsx write + read —───────────────────
    xlsx_data = {
        "sheets": {
            "E4 Results": {
                "headers": ["Model", "Baseline η", "Guarded η", "Δ"],
                "rows": [["qwen2.5:7b", 0.687, 0.952, 0.349],
                         ["mss-ai-v3.4.3", 0.771, 0.859, -0.012]],
            }
        }
    }
    xlsx_path = os.path.join(td, "e4_data.xlsx")
    r2 = agent.xlsx_write(xlsx_data, xlsx_path)
    assert r2.sheets == 1 and r2.total_rows == 2
    print(f"T3 PASS: .xlsx write ({r2.size_kb:.1f} KB, {r2.sheets} sheets)")
    
    # Read back
    r3 = agent.xlsx_read(xlsx_path)
    assert "E4 Results" in r3.sheets
    assert len(r3.sheets["E4 Results"]) == 3  # header + 2 rows
    assert r3.row_counts["E4 Results"] == 3
    print(f"T4 PASS: .xlsx read ({r3.row_counts} rows)")
    
    # ── T5: PDF generate —─────────────────────────
    md_content = """# E4 GuardProxy Summary

## Key Findings

qwen7b gains significantly under guard (+0.349), while mss-ai shows
a slight regression (-0.012) — consistent with the hypothesis that
MSS's built-in axioms conflict with external guard instructions.

## Data

| Model | Baseline | Guarded | Delta |
|-------|----------|---------|-------|
| qwen7b | 0.687 | 0.952 | +0.349 |
| mss-ai | 0.771 | 0.859 | -0.012 |
"""
    pdf_path = os.path.join(td, "e4_summary.pdf")
    r4 = agent.pdf_generate(md_content, pdf_path, title="E4 GuardProxy Summary")
    assert os.path.exists(pdf_path)
    print(f"T5 PASS: PDF generated ({r4.pages} pages, {r4.size_kb:.1f} KB)")
    
    # ── T6: PPT generate —─────────────────────────
    ppt_outline = {
        "title": "MSS-AI Experiment Pipeline",
        "subtitle": "E1-E11 Summary | June 2026",
        "slides": [
            {"title": "Identity Strength Theorem",
             "bullets": ["φ_critical = 0.660", "Crossover at ~3B params",
                        "Nested Logic Trap: η=0.913"],
             "note": "See identity_strength_theorem_v1.pdf for full paper"},
            {"title": "Social Pressure Immunity",
             "bullets": ["Waterloo paradigm confirmed", "qwen7b 0% conformity (strongest)",
                        "L2-OP eliminates mss-ai sovereignty gap"],
             "note": "E-011: 3 models × 12 questions × 3 conditions"},
            {"title": "Next Steps",
             "bullets": ["Zenodo preprint upload", "Daily context LoRA retrain",
                        "Cross-lingual anchoring validation"],
             "note": "Q3 roadmap items"},
        ]
    }
    ppt_path = os.path.join(td, "mss_summary.pptx")
    r5 = agent.ppt_generate(ppt_outline, ppt_path)
    assert os.path.exists(ppt_path)
    print(f"T6 PASS: PPT generated ({r5.slides} slides, {r5.size_kb:.1f} KB)")
    
    # ── T7: Pipeline — experiment → 3 formats ────
    pipeline_data = {
        "name": "E001_eta_calibration",
        "tables": {
            "Baseline": {
                "headers": ["Model", "η_mean", "η_std", "Breach%"],
                "rows": [["qwen7b", 1.0, 0.10, "2%"],
                         ["qwen0.5b", 0.99, 0.04, "0%"]]
            }
        },
        "report": {
            "title": "E001 Calibration",
            "author": "MSS-AI",
            "sections": [
                {"heading": "Summary", "body": "Calibration run: 10 models × 8 rounds."},
            ]
        },
        "summary_md": "# E001\n\nCalibration successful.\n"
    }
    paths = agent.experiment_to_report(pipeline_data, os.path.join(td, "pipeline_test"))
    assert len(paths) >= 2
    print(f"T7 PASS: Pipeline → {len(paths)} formats ({list(paths.keys())})")
    
    # ── T8: .xlsx empty handling —─────────────────
    r8 = agent.xlsx_read(xlsx_path)  # re-read existing
    assert "E4 Results" in r8.schema
    print(f"T8 PASS: Schema inference: {r8.schema}")
    
    print("\nS-008 Doc-Agent: all 8 tests PASSED")


if __name__ == "__main__":
    _test()
